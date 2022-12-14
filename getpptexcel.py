from typing import Mapping
from xml.sax.handler import property_interning_dict
from pptx import Presentation
import os
import pandas as pd
import numpy as np
import openpyxl
from openpyxl.styles import Font, PatternFill, Side, Alignment, Border
from openpyxl.utils import get_column_letter

writer = pd.ExcelWriter(
    r'excelinfo.xlsx', engine='openpyxl')  # 写入excel表格
colNumIdx = 0
pptIdx = 0
files = os.listdir('/Users/leslie/py/py-ppt')
files.sort()
for i in files:
    if i.endswith(".pptx"):
        pptx = Presentation(i)
        slideMap = {}  # 创建一个map对象用来存slide内容，方便后续操作
        rowNumIdx = 0
        tableIdx = 0
        for slide in pptx.slides:  # 遍历ppt页
            for shape in slide.shapes:  # 遍历页中所有方框
                cell_list = []  # 储存ppt表格的所有单元格内容
                first_cell_list = []
                if shape.has_table:  # 判断方框中是否有表格
                    if tableIdx == 0:
                        if pptIdx == 0:
                            cell_list.append("")
                            cell_list.append(i)
                        else:
                            cell_list.append(i)
                    rownum = len(shape.table.rows)  # 获取表格的行
                    colnum = len(shape.table.columns)  # 获取表格的列
                    for row in range(rownum):
                        for coloum in range(colnum):
                            if pptIdx != 0 and coloum == 0:
                                continue
                            cell = shape.table.cell(
                                row, coloum).text  # 遍历后获取单元格内容
                            cell_list.append(cell)  # 将单元格内容存到列表中
                    if pptIdx != 0:
                        if tableIdx == 0:
                            table_base = np.array(cell_list).reshape(
                                rownum+1, colnum - 1)  # 用numpy构建数组
                        else:
                            table_base = np.array(cell_list).reshape(
                                rownum, colnum - 1)  # 用numpy构建数组
                    else:
                        if tableIdx == 0:
                            table_base = np.array(cell_list).reshape(
                                rownum+1, colnum)  # 用numpy构建数组
                        else:
                            table_base = np.array(cell_list).reshape(
                                rownum, colnum)  # 用numpy构建数组
                    table = pd.DataFrame(table_base)  # 用pandas构建二维数据
                    table.to_excel(
                        writer, startrow=rowNumIdx, startcol=colNumIdx, index=False, header=False)  # 存储到表格
                    tableIdx += 1
            rowNumIdx = rowNumIdx + rownum
        if pptIdx != 0:
            colNumIdx = colNumIdx + colnum - 1
        else:
            colNumIdx = colNumIdx + colnum
        pptIdx += 1
writer.save()


def reset_color(filename):
    wb = openpyxl.load_workbook(filename)
    fill = PatternFill(  # 设置填充样式
        fill_type='solid',
        start_color='99ccff')
    diff_fill = PatternFill(  # 设置不同单元格填充样式
        fill_type='solid',
        start_color='FFC0CB')
    font = Font(size=12, bold=True)
    border = Border(top=Side(border_style='thin', color='000000'),  # 设置边框样式
                    bottom=Side(border_style='thin', color='000000'),
                    left=Side(border_style='thin', color='000000'),
                    right=Side(border_style='thin', color='000000'))
    for i in wb.sheetnames:
        ws = wb[i]
        ws.sheet_view.showGridLines = False  # 隐藏默认网线
        for i in range(1, ws.max_column+1):
            if get_column_letter(i) == 'A':
                ws.column_dimensions[get_column_letter(i)].width = 20.0
            else:
                ws.column_dimensions[get_column_letter(i)].width = 30.0
        for c in range(1, ws.max_column+1):
            for r in range(1, ws.max_row+1):  # 添加边框
                bordercell = ws.cell(r, c)
                bordercell.border = border
        for end in range(1, ws.max_column+1):
            fillcell = ws.cell(1, end)
            fillcell.fill = fill  # 填充首行
            fillcell.font = font  # 首行加粗
        for x in range(2, ws.max_row+1):
            for y in range(2, ws.max_column+1):
                if ws.cell(x, 2).value != ws.cell(x, y).value:
                    diffcell = ws.cell(x, y)
                    diffcell.fill = diff_fill
    wb.save(filename)


reset_color("excelinfo.xlsx")
