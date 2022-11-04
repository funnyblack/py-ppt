from typing import Mapping
from xml.sax.handler import property_interning_dict
from pptx import Presentation
import os
import pandas as pd
import numpy as np
import openpyxl
from openpyxl.styles import Font, PatternFill, Side, Alignment, Border
from openpyxl.utils import get_column_letter

writer = pd.ExcelWriter(
    r'excelinfo.xlsx', engine='openpyxl')  # 写入excel表格
colNumIdx = 0
pptIdx = 0
files = os.listdir('/Users/leslie/py/py-ppt')
files.sort()
model_dict = {}
for i in files:
    if i.endswith(".pptx"):
        pptx = Presentation(i)
        rowNumIdx = 1
        # 声明字典
        key_value = {}
        dict_key = ""
        for slide in pptx.slides:  # 遍历ppt页
            for shape in slide.shapes:  # 遍历页中所有方框
                cell_list = []  # 储存ppt表格的所有单元格内容
                if shape.has_table:  # 判断方框中是否有表格
                    rownum = len(shape.table.rows)  # 获取表格的行
                    colnum = len(shape.table.columns)  # 获取表格的列
                    for row in range(rownum):
                        for coloum in range(colnum):
                            if shape.table.cell(row, coloum).text == 'SPC':
                                dict_key = shape.table.cell(row, coloum+1).text
                            if pptIdx != 0 and coloum == 0:
                                continue
                            cell = shape.table.cell(
                                row, coloum).text  # 遍历后获取单元格内容
                            cell_list.append(cell)  # 将单元格内容存到列表中
                    key_value[dict_key] = cell_list
        sorted(key_value)
        titleCell = []
        if pptIdx == 0:
            model_dict = key_value
            titleCell.append("")
            titleCell.append(i)
            table_base = np.array(titleCell).reshape(
                1, 2)  # 用numpy构建数组
        else:
            titleCell.append(i)
            table_base = np.array(titleCell).reshape(
                1, 1)  # 用numpy构建数组
        table = pd.DataFrame(table_base)  # 用pandas构建二维数据
        table.to_excel(
            writer, startrow=0, startcol=colNumIdx, index=False, header=False)  # 存储到表格
        rownum = 0  # 获取表格的行
        colnum = 0  # 获取表格的列
        if pptIdx == 0:
            for key in key_value:
                rownum = int(len(key_value[key]) / 2)
                colnum = 2
                table_base = np.array(key_value[key]).reshape(
                    rownum, colnum)  # 用numpy构建数组
                table = pd.DataFrame(table_base)  # 用pandas构建二维数据
                table.to_excel(
                    writer, startrow=rowNumIdx, startcol=colNumIdx, index=False, header=False)  # 存储到表格
                rowNumIdx = rowNumIdx + rownum
        else:
            for key in model_dict:
                rownum = int(len(model_dict[key]) / 2)
                colnum = 1
                for index in key_value:
                    if key == index:
                        table_base = np.array(key_value[key]).reshape(
                            rownum, colnum)  # 用numpy构建数组
                        table = pd.DataFrame(table_base)  # 用pandas构建二维数据
                        table.to_excel(
                            writer, startrow=rowNumIdx, startcol=colNumIdx, index=False, header=False)  # 存储到表格
                        break
                rowNumIdx = rowNumIdx + rownum
        if pptIdx != 0:
            colNumIdx = colNumIdx + 1
        else:
            colNumIdx = colNumIdx + 2
        pptIdx += 1
writer.save()


def reset_color(filename):
    wb = openpyxl.load_workbook(filename)
    fill = PatternFill(  # 设置填充样式
        fill_type='solid',
        start_color='99ccff')
    diff_fill = PatternFill(  # 设置不同单元格填充样式
        fill_type='solid',
        start_color='FFC0CB')
    font = Font(size=12, bold=True)
    border = Border(top=Side(border_style='thin', color='000000'),  # 设置边框样式
                    bottom=Side(border_style='thin', color='000000'),
                    left=Side(border_style='thin', color='000000'),
                    right=Side(border_style='thin', color='000000'))
    for i in wb.sheetnames:
        ws = wb[i]
        ws.sheet_view.showGridLines = False  # 隐藏默认网线
        for i in range(1, ws.max_column+1):
            if get_column_letter(i) == 'A':
                ws.column_dimensions[get_column_letter(i)].width = 20.0
            else:
                ws.column_dimensions[get_column_letter(i)].width = 50.0
        for c in range(1, ws.max_column+1):
            for r in range(1, ws.max_row+1):  # 添加边框
                bordercell = ws.cell(r, c)
                bordercell.border = border
        for end in range(1, ws.max_column+1):
            fillcell = ws.cell(1, end)
            fillcell.fill = fill  # 填充首行
            fillcell.font = font  # 首行加粗
        for x in range(2, ws.max_row+1):
            for y in range(2, ws.max_column+1):
                if ws.cell(x, 2).value != ws.cell(x, y).value:
                    diffcell = ws.cell(x, y)
                    diffcell.fill = diff_fill
    wb.save(filename)


reset_color("excelinfo.xlsx")
